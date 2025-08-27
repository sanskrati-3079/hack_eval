# utils.py
import os
import re
import time
import shutil
import asyncio
import base64
import io
import tempfile
import glob as _glob
import subprocess
import hashlib
from typing import List, Tuple, Dict, Any, Optional

import json

from pptx import Presentation
import pypdf
from PIL import Image, ImageStat

# Optional renderers for full-page rasterization
try:
    import pypdfium2 as pdfium  # PDF page rendering
except Exception:
    pdfium = None

# Perceptual hash for dedup (used in other modules; safe to keep here)
try:
    import imagehash  # type: ignore
except Exception:
    imagehash = None

# Optional COM for Windows PowerPoint export
try:
    import comtypes.client as _com_client  # Windows PowerPoint slide export
except Exception:
    _com_client = None


# ===== Evaluation schema =====
EVAL_WEIGHTS = {
    'Problem Understanding': 15,
    'Innovation & Uniqueness': 20,
    'Technical Feasibility': 20,
    'Implementation Approach': 15,
    'Team Readiness': 15,
    'Potential Impact': 15
}
EVAL_ORDER = [
    'Innovation & Uniqueness',
    'Technical Feasibility',
    'Potential Impact',
    'Problem Understanding',
    'Implementation Approach',
    'Team Readiness',
]

ALLOWED_EXTS = {".pdf", ".pptx", ".ppt"}

def raw_total(scores: dict) -> float:
    return round(sum(float(scores.get(k, 0)) for k in EVAL_WEIGHTS.keys()), 1)

def weighted_total(scores: dict) -> float:
    total = 0.0
    for k, w in EVAL_WEIGHTS.items():
        v = float(scores.get(k, 0))
        v = max(1.0, min(10.0, v)) if scores.get(k, None) is not None else 0.0
        total += (v / 10.0) * float(w)
    return round(total, 2)

def _stable_small_jitter(team_name: str) -> float:
    # 0..0.009 deterministic jitter to break deep ties
    h = hashlib.sha256(team_name.encode("utf-8")).hexdigest()
    val = int(h[:6], 16) / float(0xFFFFFF)  # 0..1
    return round(val * 0.009, 6)

def tie_break_key(scores: dict, team_name: str):
    def s(k):
        v = float(scores.get(k, 0))
        return max(0, min(10, v))
    return (
        -weighted_total(scores),
        -s('Innovation & Uniqueness'),
        -s('Technical Feasibility'),
        -s('Potential Impact'),
        -s('Problem Understanding'),
        -s('Implementation Approach'),
        -s('Team Readiness'),
        -_stable_small_jitter(team_name),
        team_name.lower(),
    )

def _score_total(scores: dict) -> float:
    return weighted_total(scores)


# ---------- Heuristic baselines and calibration ----------
_KEYWORDS = {
    'Problem Understanding': ["problem", "challenge", "pain point", "scope", "objective", "goal", "requirement"],
    'Innovation & Uniqueness': ["novel", "unique", "innovation", "innovative", "differentiator", "patent", "state-of-the-art", "sota", "first-of-its-kind"],
    'Technical Feasibility': ["architecture", "api", "dataset", "model", "deployment", "latency", "scalability", "prototype", "poc", "throughput", "memory", "inference", "kubernetes", "k8s", "terraform", "etl"],
    'Implementation Approach': ["timeline", "milestone", "component", "module", "pipeline", "workflow", "algorithm", "stack", "roadmap", "sprint"],
    'Team Readiness': ["team", "member", "experience", "skills", "roles", "responsibilities", "domain expertise"],
    'Potential Impact': ["impact", "roi", "revenue", "users", "adoption", "market", "metrics", "kpi", "benchmark", "evaluation", "a/b", "ab test", "growth", "auc", "bleu"],
}
_EXTRA_EVIDENCE = ["baseline", "privacy", "security", "gdpr", "hipaa", "cost", "budget", "infra", "cloud", "risk", "mitigation"]

def _contains_any(text_lc: str, words: List[str]) -> bool:
    return any(w.lower() in text_lc for w in words)

def _count_numbers(text: str) -> int:
    return len(re.findall(r"\b\d+(?:\.\d+)?%?\b", text or ""))

def _technical_density(text: str) -> float:
    words = re.findall(r"[a-zA-Z0-9\-]+", text or "")
    if not words:
        return 0.0
    tech_hits = sum(1 for w in words if w.lower() in {
        "api","kpi","roc","auc","bleu","etl","k8s","kubernetes","terraform","latency","throughput","inference","model","dataset"
    })
    return tech_hits / len(words)

def _heuristic_baseline(raw_text: str, images_count: int) -> dict:
    text_lc = (raw_text or "").lower()
    wc = len((raw_text or "").split())
    nums = _count_numbers(raw_text)
    density = _technical_density(raw_text)

    if wc >= 400:
        base = 6
    elif wc >= 200:
        base = 5
    elif wc >= 100:
        base = 4
    else:
        base = 3

    if density < 0.01 and wc > 150:
        base = max(3, base - 1)

    adj_img = 1 if images_count > 0 else 0

    baseline = {}
    for k, kws in _KEYWORDS.items():
        score = base
        if _contains_any(text_lc, kws):
            score += 1
        if k in ("Technical Feasibility", "Potential Impact") and nums >= 2:
            score += 1
        if k in ("Technical Feasibility", "Implementation Approach"):
            score += adj_img
        if k == "Innovation & Uniqueness" and _contains_any(text_lc, ["novel", "unique", "patent", "state-of-the-art", "sota", "first"]):
            score += 1
        baseline[k] = int(max(3, min(8, score)))
    return baseline

def calibrate_and_enrich_scores(raw_text: str, images_count: int, scores: dict) -> dict:
    text = raw_text or ""
    text_lc = text.lower()
    wc = len(text.split())
    nums = _count_numbers(text)
    density = _technical_density(text)
    baseline = _heuristic_baseline(text, images_count)

    s = {}
    for k in EVAL_WEIGHTS.keys():
        v = scores.get(k, None)
        if v is None or not isinstance(v, (int, float)) or v < 1 or v > 10:
            s[k] = baseline[k]
        else:
            s[k] = float(v)

    cap = 9 if wc < 150 else 10
    for k in s:
        s[k] = min(s[k], cap)

    if not _contains_any(text_lc, _EXTRA_EVIDENCE):
        for k in ["Technical Feasibility", "Potential Impact"]:
            s[k] = max(3.0, s[k] - 1.0)

    if nums >= 5:
        for k in ["Potential Impact", "Technical Feasibility"]:
            s[k] = min(10.0, s[k] + 1.0)

    if density < 0.01 and wc > 150:
        for k in ["Technical Feasibility", "Implementation Approach"]:
            s[k] = max(3.0, s[k] - 1.0)

    if wc >= 120 or images_count > 0:
        for k in s:
            s[k] = max(3.0, s[k])

    tens = [k for k, v in s.items() if round(v) >= 10]
    if len(tens) > 1:
        priority = ['Innovation & Uniqueness', 'Technical Feasibility', 'Potential Impact',
                    'Problem Understanding', 'Implementation Approach', 'Team Readiness']
        keep = next((p for p in priority if p in tens), tens[0])
        for k in tens:
            if k != keep:
                s[k] = 9.0

    uniq_vals = {int(round(v)) for v in s.values()}
    if len(uniq_vals) == 1:
        for k in ['Team Readiness', 'Implementation Approach', 'Problem Understanding']:
            s[k] = max(3.0, s[k] - 1.0)

    out = {k: int(min(10, max(1, round(v)))) for k, v in s.items()}
    return out


# ---------- Rate limiters: text vs vision ----------
class _RateLimiter:
    def __init__(self, rpm: int):
        self.min_interval = 60.0 / max(1, rpm)
        self._last_ts = 0.0
        self._lock = asyncio.Lock()

    async def acquire(self):
        async with self._lock:
            now = time.perf_counter()
            wait = self._last_ts + self.min_interval - now
            if wait > 0:
                await asyncio.sleep(wait)
            self._last_ts = time.perf_counter()

_TEXT_LIMITER = _RateLimiter(int(os.getenv("RATE_LIMIT_RPM_TEXT", "18")))
_VISION_LIMITER = _RateLimiter(int(os.getenv("RATE_LIMIT_RPM_VISION", "6")))

def get_text_limiter():
    return _TEXT_LIMITER

def get_vision_limiter():
    return _VISION_LIMITER


# ---------- JSON extraction helper ----------
def extract_first_json_object(text: str) -> Optional[str]:
    """Return the first top-level JSON object found in text, or None."""
    if not text:
        return None
    # Prefer fenced blocks first
    if "```json" in text:
        try:
            chunk = text.split("```json", 1)[1].split("```", 1)[0].strip()
            if chunk.startswith("{") and chunk.endswith("}"):
                return chunk
        except Exception:
            pass

    # Brace-matching scan that respects strings and escapes
    start = -1
    depth = 0
    in_str = False
    escape = False
    for i, ch in enumerate(text):
        if start == -1:
            if ch == '{':
                start = i
                depth = 1
                in_str = False
                escape = False
            continue

        if escape:
            escape = False
            continue

        if ch == '\\':
            if in_str:
                escape = True
            continue

        if ch == '"':
            in_str = not in_str
            continue

        if not in_str:
            if ch == '{':
                depth += 1
            elif ch == '}':
                depth -= 1
                if depth == 0:
                    return text[start:i+1]

    return None


# ---------- Image utilities ----------
def _phash(pil: Image.Image) -> Optional[str]:
    try:
        if imagehash is None:
            return None
        return str(imagehash.phash(pil))
    except Exception:
        return None

def _is_decorative(pil: Image.Image) -> bool:
    # Heuristic: very small, very low-contrast, or almost-solid images are likely decorative
    w, h = pil.size
    if w * h < 30000:  # tiny
        return True
    stat = ImageStat.Stat(pil.convert("L"))
    if stat.var[0] < 50:  # very low variance
        return True
    return False

def _to_b64_jpeg(pil_img: Image.Image, quality=85) -> str:
    if pil_img.mode == "RGBA":
        pil_img = pil_img.convert("RGB")
    buf = io.BytesIO()
    pil_img.save(buf, format="JPEG", quality=quality)
    return base64.b64encode(buf.getvalue()).decode("utf-8")


# ---------- PDF loaders ----------
def _render_pdf_pages_to_images(path: str, dpi: int = 150) -> List[Dict[str, Any]]:
    """Render each PDF page to an image dict with base64 and page_index."""
    if pdfium is None:
        return []
    out: List[Dict[str, Any]] = []
    pdf = pdfium.PdfDocument(path)
    for i in range(len(pdf)):
        page = pdf[i]
        pil = page.render(scale=dpi / 72.0).to_pil().convert("RGB")
        if not _is_decorative(pil):
            out.append({"b64": _to_b64_jpeg(pil, quality=85), "page_index": i})
    return out

def _extract_pdf_text_and_images(path: str) -> Tuple[str, List[str]]:
    """Text + visuals for evidence count. Always include page renders to capture vector diagrams."""
    text_parts: List[str] = []
    images_b64: List[str] = []
    with open(path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for page in reader.pages:
            try:
                text_parts.append(page.extract_text() or "")
            except Exception as e:
                print(f"[pdf text warn] {type(e).__name__}")
        # Embedded raster images
        try:
            for page in reader.pages:
                if "/Resources" in page and "/XObject" in page["/Resources"]: # pyright: ignore[reportOperatorIssue]
                    xobj = page["/Resources"]["/XObject"].get_object() # pyright: ignore[reportIndexIssue]
                    for obj in xobj:
                        o = xobj[obj]
                        if o.get("/Subtype") == "/Image":
                            data = o.get_data()
                            pil = Image.open(io.BytesIO(data)).convert("RGB")
                            if not _is_decorative(pil):
                                images_b64.append(_to_b64_jpeg(pil))
        except Exception as e:
            print(f"[pdf img warn] {type(e).__name__}")
    # Always render pages as well (captures SmartArt/vector)
    pages = _render_pdf_pages_to_images(path)
    max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
    if max_pages > 0:
        pages = pages[:max_pages]
    images_b64.extend([p["b64"] for p in pages])
    return "\n".join(text_parts), images_b64


# ---------- PPT/PPTX loaders ----------
def _render_pptx_slides_windows(path: str) -> List[Dict[str, Any]]:
    """Windows PowerPoint COM export. Clean up reliably."""
    if _com_client is None or os.name != "nt":
        return []
    tmpdir = tempfile.mkdtemp(prefix="ppt_render_")
    out: List[Dict[str, Any]] = []
    pp = None
    pres = None
    try:
        pp = _com_client.CreateObject("PowerPoint.Application")
        pp.Visible = 0
        pres = pp.Presentations.Open(path, WithWindow=False)
        pres.Export(tmpdir, "PNG")
        # Gather PNGs in slide order
        files = sorted(_glob.glob(os.path.join(tmpdir, "*.PNG")) + _glob.glob(os.path.join(tmpdir, "*.png")))
        for idx, png in enumerate(files):
            try:
                pil = Image.open(png).convert("RGB")
                if not _is_decorative(pil):
                    out.append({"b64": _to_b64_jpeg(pil), "slide_index": idx})
            except Exception as e:
                print(f"[pptx render warn] slide {idx}: {type(e).__name__}")
    except Exception as e:
        print(f"[pptx render error] {type(e).__name__}: {e}")
    finally:
        try:
            if pres: pres.Close()
        except Exception:
            pass
        try:
            if pp: pp.Quit()
        except Exception:
            pass
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out

def _render_with_soffice(path: str) -> List[Dict[str, Any]]:
    """LibreOffice headless export for PPT/PPTX to PNGs."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return []
    tmpdir = tempfile.mkdtemp(prefix="soffice_")
    out: List[Dict[str, Any]] = []
    try:
        # Convert to PNGs
        cmd = [soffice, "--headless", "--convert-to", "png", "--outdir", tmpdir, path]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        files = sorted(_glob.glob(os.path.join(tmpdir, "*.png")))
        for idx, png in enumerate(files):
            try:
                pil = Image.open(png).convert("RGB")
                if not _is_decorative(pil):
                    out.append({"b64": _to_b64_jpeg(pil), "slide_index": idx})
            except Exception as e:
                print(f"[soffice warn] slide {idx}: {type(e).__name__}")
    except Exception as e:
        print(f"[soffice error] {type(e).__name__}: {e}")
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)
    return out

def _extract_pptx_text_and_images(path: str) -> Tuple[str, List[str]]:
    """Text + visuals for evidence count. Always add rendered slides to capture shapes/SmartArt."""
    text_parts: List[str] = []
    images_b64: List[str] = []
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"[pptx open warn] {type(e).__name__}")
        prs = None
    if prs:
        for s_i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and getattr(shape, "has_text_frame", False):
                    try:
                        text_parts.append(shape.text) # pyright: ignore[reportAttributeAccessIssue]
                    except Exception:
                        pass
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    try:
                        pil = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                        if not _is_decorative(pil):
                            images_b64.append(_to_b64_jpeg(pil))
                    except Exception:
                        continue
    # Always render slides too
    rendered = _render_pptx_slides_windows(path) if os.name == "nt" else _render_with_soffice(path)
    max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
    if max_pages > 0:
        rendered = rendered[:max_pages]
    images_b64.extend([r["b64"] for r in rendered])
    return "\n".join(text_parts), images_b64

def load_document_content(file_path: str) -> Tuple[str, List[str]]:
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in ALLOWED_EXTS:
        return "", []
    if ext == ".pdf":
        return _extract_pdf_text_and_images(file_path)
    if ext in (".pptx", ".ppt"):
        return _extract_pptx_text_and_images(file_path)
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []
    except Exception:
        return "", []


# ---------- Reporting ----------
def display_consolidated_report(context) -> None:
    print(f"\nTeam: {context.team_name}")
    if context.evaluation_error:
        print(f"Error: {context.evaluation_error}\n")
        return

    print("Scores:")
    for k in EVAL_WEIGHTS.keys():
        v = context.scores.get(k, "-")
        print(f"  - {k}: {v}")
    print(f"Total (raw): {raw_total(context.scores)}")
    print(f"Total (weighted): {weighted_total(context.scores)}")
    print(f"Summary: {context.scoring_summary or ''}")

    overall = None
    if isinstance(context.workflow_analysis, dict):
        overall = context.workflow_analysis.get("overall_summary")
    if not overall and context.workflow_report:
        overall = context.workflow_report.get("overall_summary")
    if not overall:
        overall = "Null. No diagrams were provided in the presentation text."
    print("Workflow Analysis:")
    print(f"  Overall: {overall}")

    if context.feedback:
        print("Feedback:")
        print(f"  Positive: {context.feedback.get('positive','')}")
        print(f"  Criticism: {context.feedback.get('criticism','')}")
        print(f"  Technical: {context.feedback.get('technical','')}")
        print(f"  Suggestions: {context.feedback.get('suggestions','')}")
    print()

    # JSON output
    json_report = {
        "team_name": context.team_name,
        "scores": {k: context.scores.get(k, "-") for k in EVAL_WEIGHTS.keys()},
        "total_raw": raw_total(context.scores),
        "total_weighted": weighted_total(context.scores),
        "summary": context.scoring_summary or "",
        "workflow_analysis": {
            "overall": overall
        },
        "feedback": context.feedback or {}
    }
    print("JSON format:")
    print(json.dumps(json_report, indent=2))
    
def display_leaderboard(contexts: list) -> None:
    printable = []
    for c in contexts:
        total = _score_total(c.scores) if not c.evaluation_error else -1
        printable.append((total, c.team_name, c.file_path, c.scores or {}))
    printable.sort(key=lambda x: tie_break_key(x[3], x[1]))

    print("\n######## Leaderboard ########")
    for i, (total, name, path, scores) in enumerate(printable, 1):
        status = f"{total:.2f}" if total >= 0 else "ERROR"
        inz = scores.get('Innovation & Uniqueness', '-')
        tech = scores.get('Technical Feasibility', '-')
        imp = scores.get('Potential Impact', '-')
        print(f"{i:2d}. {name:30s} | weighted {status:>6s} | Inno {inz} | Tech {tech} | Impact {imp} | {os.path.basename(path)}")
    print("#############################\n")


# ...existing code...

def save_consolidated_reports_to_excel(contexts: list, filename: str):
    """Save all consolidated report data to an Excel file."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not installed. Skipping Excel export.")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Reports"

    # Header
    headers = [
        "team_name", "file_path", "evaluation_error"
    ] + list(EVAL_WEIGHTS.keys()) + [
        "total_raw", "total_weighted", "summary", "workflow_overall",
        "feedback_positive", "feedback_criticism", "feedback_technical", "feedback_suggestions"
    ]
    ws.append(headers)

    for ctx in contexts:
        row = [
            getattr(ctx, "team_name", ""),
            getattr(ctx, "file_path", ""),
            getattr(ctx, "evaluation_error", ""),
        ]
        for k in EVAL_WEIGHTS.keys():
            row.append(ctx.scores.get(k, "-") if ctx.scores else "-")
        row.append(raw_total(ctx.scores) if ctx.scores else "-")
        row.append(weighted_total(ctx.scores) if ctx.scores else "-")
        row.append(getattr(ctx, "scoring_summary", ""))
        # Workflow overall summary
        overall = None
        if isinstance(ctx.workflow_analysis, dict):
            overall = ctx.workflow_analysis.get("overall_summary")
        if not overall and getattr(ctx, "workflow_report", None):
            overall = ctx.workflow_report.get("overall_summary")
        if not overall:
            overall = "Null. No diagrams were provided in the presentation text."
        row.append(overall)
        # Feedback
        fb = ctx.feedback or {}
        row.append(fb.get("positive", ""))
        row.append(fb.get("criticism", ""))
        row.append(fb.get("technical", ""))
        row.append(fb.get("suggestions", ""))
        ws.append(row)

    wb.save(filename)
    print(f"[Excel] Consolidated reports saved to {filename}")

def save_leaderboard_to_excel(contexts: list, filename: str):
    """Save leaderboard data to an Excel file."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not installed. Skipping Excel export.")
        return

    wb = Workbook()
    ws = wb.active
    ws.title = "Leaderboard"

    headers = [
        "Rank", "Team Name", "Weighted Total", "Innovation & Uniqueness",
        "Technical Feasibility", "Potential Impact", "File Name"
    ]
    ws.append(headers)

    printable = []
    for c in contexts:
        total = _score_total(c.scores) if not c.evaluation_error else -1
        printable.append((total, c.team_name, c.file_path, c.scores or {}))
    printable.sort(key=lambda x: tie_break_key(x[3], x[1]))

    for i, (total, name, path, scores) in enumerate(printable, 1):
        status = f"{total:.2f}" if total >= 0 else "ERROR"
        inz = scores.get('Innovation & Uniqueness', '-')
        tech = scores.get('Technical Feasibility', '-')
        imp = scores.get('Potential Impact', '-')
        ws.append([i, name, status, inz, tech, imp, os.path.basename(path)])

    wb.save(filename)
    print(f"[Excel] Leaderboard saved to {filename}")


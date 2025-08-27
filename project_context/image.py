import os
import base64
import io
from dotenv import load_dotenv
from pydantic.v1 import BaseModel, Field
from typing import List, Dict
from pptx import Presentation
import pypdf
from PIL import Image

from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.messages import HumanMessage

load_dotenv()

# ---------------------- SCHEMA ----------------------

class ImageAnalysis(BaseModel):
    file: str = Field(description="File name or reference of the image.")
    diagram_type: str = Field(description="Diagram type, e.g., Flowchart, Architecture, User Flow.")
    step_by_step: List[str] = Field(description="Step-by-step explanation of the process.")
    observations: Dict[str, List[str]] = Field(description="Grouped into positive, criticism, technical.")
    suggestions: List[str] = Field(description="Bold, short, direct suggestions.")

class WorkflowReport(BaseModel):
    images: List[ImageAnalysis] = Field(description="Per-image analyses.", min_items=1)
    overall_summary: Dict[str, List[str]] = Field(
        description="Overall summary grouped into positive, criticism, technical, suggestions."
    )

# ---------------------- AGENT ----------------------

class WorkflowAnalysisAgent:
    def __init__(self):
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY not found in .env file.")

        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.2,
            openai_api_key=api_key,
        )
        self.parser = JsonOutputParser(pydantic_object=WorkflowReport)
        self.prompt = self._create_prompt()

    def _create_prompt(self):
        prompt_str = """
You are a **System Design and Process Analysis Specialist**.  
You analyze images extracted from a presentation (flowcharts, architecture diagrams, user journeys, etc.).  

Your output must follow these rules:

1. **Image Analysis (per image)**  
   - Classify the diagram type (**Flowchart**, **Architecture**, **User Journey**, etc.).  
   - Explain the process in a **step-by-step list** (each step short & simple).  
   - Provide observations in 3 groups:  
     - **Positive** → good practices.  
     - **Criticism** → missing parts or weaknesses.  
     - **Technical** → system/architecture details or gaps.  
   - Provide **Suggestions**:  
     - Each suggestion must be **bold**, **short**, and **direct**.  
     - If multiple, list them point-wise.  

2. **Overall Summary**  
   - Combine insights into 4 categories: **Positive**, **Criticism**, **Technical**, **Suggestions**.  
   - Keep bullet-pointed and brief.  
   - Suggestions must be **bold and actionable**.  

3. **Formatting Rules**  
   - Use **bold titles** for sections.  
   - Use bullet points for all lists.  
   - Keep sentences **short and simple**.  

Return ONLY JSON matching this schema:
{format_instructions}
"""
        return ChatPromptTemplate.from_template(prompt_str)

    def _extract_images_as_base64(self, file_path):
        images = []
        print(f"  -> Extracting images from '{file_path}'...")
        try:
            if file_path.lower().endswith(".pdf"):
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    if hasattr(page, "images") and page.images:
                        for image_file_object in page.images:
                            img = Image.open(io.BytesIO(image_file_object.data))
                            if img.mode == "RGBA":
                                img = img.convert("RGB")
                            buffered = io.BytesIO()
                            img.save(buffered, format="JPEG")
                            images.append(base64.b64encode(buffered.getvalue()).decode("utf-8"))

            elif file_path.lower().endswith(".pptx"):
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            image_bytes = shape.image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            if img.mode == "RGBA":
                                img = img.convert("RGB")
                            buffered = io.BytesIO()
                            img.save(buffered, format="JPEG")
                            images.append(base64.b64encode(buffered.getvalue()).decode("utf-8"))
        except Exception as e:
            print(f"  -> Warning: Could not extract images. {e}")

        print(f"  -> Found {len(images)} images.")
        return images

    def analyze_workflows(self, file_path):
        images_base64 = self._extract_images_as_base64(file_path)
        if not images_base64:
            print("  -> No images found to analyze.")
            return None

        prompt_text = self.prompt.format(format_instructions=self.parser.get_format_instructions())
        message_parts = [{"type": "text", "text": prompt_text}]
        for img_data in images_base64:
            message_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{img_data}"}
            })

        message = HumanMessage(content=message_parts)

        print("  -> Calling OpenAI API for workflow analysis...")
        try:
            response = self.llm.invoke([message])
            raw_content = response.content

            if "```json" in raw_content:
                clean_content = raw_content.split("```json")[1].split("```")[0].strip()
            else:
                clean_content = raw_content

            report_data = self.parser.parse(clean_content)
            print("  -> Analysis complete.")
            return WorkflowReport(**report_data)

        except Exception as e:
            print(f"  -> ERROR during workflow analysis: {e}")
            return None

# ---------------------- DISPLAY ----------------------

def display_workflow_report(report: WorkflowReport):
    if not report:
        print("Could not generate a workflow report.")
        return

    print("\n" + "=" * 70)
    print("📊 WORKFLOW AND DIAGRAM ANALYSIS REPORT 📊")
    print("=" * 70)

    print("\n--- 📜 Overall Workflow Summary ---")
    for section, points in report.overall_summary.items():
        print(f"\n**{section.capitalize()}**")
        for p in points:
            print(f"- {p}")

    print("\n--- 🖼️ Detailed Image Analysis ---")
    for analysis in report.images:
        print(f"\n➡️ Image ({analysis.diagram_type}):")
        print("Step by Step:")
        for step in analysis.step_by_step:
            print(f"- {step}")
        print("\nObservations:")
        for sec, pts in analysis.observations.items():
            print(f"  **{sec.capitalize()}**")
            for p in pts:
                print(f"  - {p}")
        print("\nSuggestions:")
        for s in analysis.suggestions:
            print(f"- {s}")

    print("\n" + "=" * 70)


if __name__ == "__main__":
    TEAM_FILE_PATH = r"C:\Users\risha\Downloads\project_context\ppt\Algoknights_Ideaformat[1][1] - SARANSH GUPTA.pdf"

    print(f"--- Starting Workflow Analysis for: {TEAM_FILE_PATH} ---")

    if not os.path.exists(TEAM_FILE_PATH):
        print(f"FATAL ERROR: The file '{TEAM_FILE_PATH}' was not found.")
    else:
        agent = WorkflowAnalysisAgent()
        workflow_report = agent.analyze_workflows(TEAM_FILE_PATH)
        if workflow_report:
            display_workflow_report(workflow_report)

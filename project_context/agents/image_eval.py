import os
import io
import base64
from typing import List, Dict, Any, Optional

from pptx import Presentation
import pypdf
from PIL import Image
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.messages import HumanMessage
from pydantic.v1 import BaseModel, Field

# Optional helpers
try:
    import pypdfium2 as pdfium # pyright: ignore[reportMissingImports]
except Exception:
    pdfium = None

try:
    import imagehash  # type: ignore
except Exception:
    imagehash = None

def _to_b64_jpeg(pil_img: Image.Image, quality=85) -> str:
    if pil_img.mode == "RGBA":
        pil_img = pil_img.convert("RGB")
    buf = io.BytesIO()
    pil_img.save(buf, format="JPEG", quality=quality)
    return base64.b64encode(buf.getvalue()).decode("utf-8")

def _phash(pil: Image.Image) -> Optional[str]:
    try:
        if imagehash is None:
            return None
        return str(imagehash.phash(pil))
    except Exception:
        return None

class ImageAnalysis(BaseModel):
    image_index: int = Field(description="Index number of the image being analyzed.")
    description: str = Field(description="Step-by-step description of the diagram or image.")
    type: str = Field(description="Diagram type, e.g., Architecture, User Flow, Data Flow, Sequence, Chart, Mockup.")
    slide_index: Optional[int] = Field(default=None, description="0-based slide index if from PPT/PPTX.")
    page_index: Optional[int] = Field(default=None, description="0-based page index if from PDF.")
    is_diagram: bool = Field(description="True if the image is a diagram/flow/architecture, False if photo/logo.")
    importance: str = Field(description="One of: critical, supporting, decorative, irrelevant.")
    confidence: float = Field(description="0.0–1.0 confidence in the classification.", ge=0.0, le=1.0)

class WorkflowReport(BaseModel):
    overall_summary: str = Field(description="High-level workflow across all diagrams.")
    image_analyses: List[ImageAnalysis] = Field(description="Per-image analyses with classification.", min_items=1)

class WorkflowAnalysisAgent:
    """
    Captures ALL deck visuals:
    - Always combines embedded images + rendered full slides/pages.
    - Dedups via perceptual hash and keeps slide order.
    - Classifies diagram vs photo and importance.
    """
    def __init__(self):
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY not found in .env file.")
        self.vision_model = os.getenv("OPENAI_MODEL_VISION", os.getenv("OPENAI_MODEL", "gpt-4o"))
        self.llm = ChatOpenAI(
            model=self.vision_model,
            temperature=0.2,
            top_p=0.0,
            api_key=api_key, # pyright: ignore[reportArgumentType]
        )
        self.parser = JsonOutputParser(pydantic_object=WorkflowReport)
        self.prompt = self._create_prompt()
        self.max_images = int(os.getenv("MAX_VISION_IMAGES", "12"))

    def _create_prompt(self):
        prompt_str = """
You are a system design and process analysis specialist.
Analyze the images extracted from a presentation. These images can be flowcharts, architecture diagrams,
user journeys, data/ML pipelines, deployment diagrams, charts, mockups, or photos.

Instructions:
- First classify each image: is_diagram = true if it has boxes/arrows/lanes/data flows/components;
  false if it is a photo, logo, or decorative picture.
- Set importance:
  - "critical" if it captures the core system/workflow,
  - "supporting" if it explains a component or subflow,
  - "decorative" for screenshots/illustrations with little process information,
  - "irrelevant" if it does not relate to the project.
- Provide a step-by-step description of each image that is a diagram. Summaries for photos should be brief.
- Include slide/page indices when present.
- After all analyses, produce an overall workflow summary that relies mainly on critical/supporting diagrams.
- Give confidence 0.0–1.0 for your classification.

Return ONLY JSON matching this schema:

{format_instructions}
"""
        return ChatPromptTemplate.from_template(prompt_str)

    # -------- PDF helpers --------
    def _extract_pdf_embedded(self, file_path: str) -> List[Dict[str, Any]]:
        out = []
        try:
            reader = pypdf.PdfReader(file_path)
            for p_i, page in enumerate(reader.pages):
                try:
                    if "/Resources" in page and "/XObject" in page["/Resources"]:
                        xobj = page["/Resources"]["/XObject"].get_object()
                        for obj in xobj:
                            o = xobj[obj]
                            if o.get("/Subtype") == "/Image":
                                data = o.get_data()
                                pil = Image.open(io.BytesIO(data)).convert("RGB")
                                out.append({"b64": _to_b64_jpeg(pil), "page_index": p_i, "ph": _phash(pil)})
                except Exception:
                    continue
        except Exception as e:
            print(f"[img pdf embedded warn] {type(e).__name__}: {e}")
        return out

    def _render_pdf_pages(self, file_path: str, dpi: int = 170) -> List[Dict[str, Any]]:
        if pdfium is None:
            return []
        out = []
        try:
            pdf = pdfium.PdfDocument(file_path)
            max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
            for i in range(min(len(pdf), max_pages)):
                page = pdf[i]
                pil = page.render(scale=dpi / 72.0).to_pil().convert("RGB")
                out.append({"b64": _to_b64_jpeg(pil), "page_index": i, "ph": _phash(pil)})
        except Exception as e:
            print(f"[img pdf render warn] {type(e).__name__}: {e}")
        return out

    # -------- PPTX helpers --------
    def _extract_pptx_embedded(self, file_path: str) -> List[Dict[str, Any]]:
        out = []
        try:
            prs = Presentation(file_path)
            for s_i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            pil = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                            out.append({"b64": _to_b64_jpeg(pil), "slide_index": s_i, "ph": _phash(pil)})
                        except Exception:
                            continue
        except Exception as e:
            print(f"[img pptx embedded warn] {type(e).__name__}: {e}")
        return out

    def _render_ppt_generic(self, file_path: str) -> List[Dict[str, Any]]:
        # Windows: COM via utils; Other OS: LibreOffice headless
        out = []
        if os.name == "nt":
            try:
                from utils import _render_pptx_slides_windows as _render_win
                out = _render_win(file_path)
                for d in out:
                    d["ph"] = None
            except Exception as e:
                print(f"[img pptx COM fallback warn] {type(e).__name__}: {e}")
        else:
            try:
                from utils import _render_with_soffice as _render_soffice
                out = _render_soffice(file_path)
                for d in out:
                    d["ph"] = None
            except Exception as e:
                print(f"[img soffice fallback warn] {type(e).__name__}: {e}")
        # limit pages
        max_pages = int(os.getenv("MAX_RENDER_PAGES", "12"))
        return out[:max_pages] if max_pages > 0 else out

    # -------- merge + dedup --------
    def _dedup_and_order(self, images: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        seen = set()
        uniq = []
        for d in images:
            ph = d.get("ph")
            key = ph or (d.get("slide_index"), d.get("page_index"), len(d.get("b64","")))
            if key in seen:
                continue
            seen.add(key)
            uniq.append(d)
        # Keep rendered slide/page images first so diagrams survive logos
        uniq.sort(key=lambda x: (
            0 if ("slide_index" in x or "page_index" in x) else 1,  # rendered first
            x.get("slide_index", 1_000_000),
            x.get("page_index", 1_000_000),
        ))
        if len(uniq) > self.max_images:
            step = max(1, len(uniq) // self.max_images)
            uniq = uniq[::step][:self.max_images]
        return uniq

    def _extract_images_as_base64(self, file_path: str) -> List[Dict[str, Any]]:
        print(f"  -> Extracting images from '{file_path}'...")
        lower = file_path.lower()

        embedded: List[Dict[str, Any]] = []
        rendered: List[Dict[str, Any]] = []

        try:
            if lower.endswith(".pdf"):
                embedded = self._extract_pdf_embedded(file_path)
                rendered = self._render_pdf_pages(file_path)
            elif lower.endswith(".pptx") or lower.endswith(".ppt"):
                embedded = self._extract_pptx_embedded(file_path)
                rendered = self._render_ppt_generic(file_path)
        except Exception as e:
            print(f"  -> Warning: Could not extract images. {e}")

        combined = rendered + embedded  # rendered first to prioritize diagrams
        images = self._dedup_and_order(combined)
        print(f"  -> Using {len(images)} image(s) for analysis.")
        return images

    # -------- LLM call --------
    def analyze_workflows(self, file_path: str) -> Optional[WorkflowReport]:
        images = self._extract_images_as_base64(file_path)
        if not images:
            print("  -> No images found to analyze.")
            return None

        prompt_text = self.prompt.format(format_instructions=self.parser.get_format_instructions())
        parts: List[Dict[str, Any]] = [{"type": "text", "text": prompt_text}]
        for idx, d in enumerate(images, start=1):
            parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{d['b64']}", "detail": "high"}
            })
            meta_bits = []
            if d.get("slide_index") is not None:
                meta_bits.append(f"(slide {d['slide_index']+1})")
            if d.get("page_index") is not None:
                meta_bits.append(f"(page {d['page_index']+1})")
            if meta_bits:
                parts.append({"type": "text", "text": f"Image {idx} context: {' '.join(meta_bits)}"})

        message = HumanMessage(content=parts)

        print("  -> Calling OpenAI API for workflow analysis...")
        try:
            resp = self.llm.invoke([message])
            raw = resp.content or ""
            clean = raw
            if "```json" in raw:
                try:
                    clean = raw.split("```json", 1)[1].split("```", 1)[0].strip()
                except Exception:
                    clean = raw

            data = self.parser.parse(clean)

            # Attach indices + defaults
            enriched = []
            for i, ia in enumerate(data.get("image_analyses", []), start=1):
                meta = images[i-1] if i-1 < len(images) else {}
                ia.setdefault("slide_index", meta.get("slide_index"))
                ia.setdefault("page_index", meta.get("page_index"))
                ia.setdefault("is_diagram", (ia.get("type","").lower() not in ("photo","image","mockup")))
                ia.setdefault("importance", "supporting" if ia["is_diagram"] else "decorative")
                ia.setdefault("confidence", 0.7)
                enriched.append(ia)
            data["image_analyses"] = enriched

            print("  -> Analysis complete.")
            return WorkflowReport(**data)

        except Exception as e:
            print(f"  -> ERROR during workflow analysis: {type(e).__name__}: {e}")
            return None
